from pptx import Presentation

prs = Presentation()

# 標題投影片
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
title.text = "個人履歷"
subtitle = slide1.placeholders[1]
subtitle.text = "王小明 - Python 開發者"

# 內容投影片
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide2.shapes.title
title.text = "個人資訊"
content = slide2.placeholders[1]
content.text = "電話: 0987-654-321\nEmail: example@email.com\n技能: Python, AI, 數據分析"

prs.save("個人履歷簡報.pptx")
print("已建立個人履歷簡報.pptx")
