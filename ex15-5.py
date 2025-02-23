from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

prs = Presentation("簡報.pptx")
slide = prs.slides.add_slide(prs.slide_layouts[5])  # 空白投影片
chart_data = CategoryChartData()
chart_data.categories = ["A", "B", "C"]
chart_data.add_series("數據", (30, 50, 40))

x, y, cx, cy = 100, 100, 400, 300  # 設定圖表位置
slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

prs.save("簡報_圖表.pptx")
print("已新增圖表至 PowerPoint")
